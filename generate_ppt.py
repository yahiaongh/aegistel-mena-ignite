from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ==============================================================================
# DESIGN SYSTEM — Restrained, editorial, whitespace-first
# ==============================================================================

# Palette: One background, one accent, three text values.
C_BG      = RGBColor(8, 14, 26)      # Deep space navy — calmer than pure black
C_ACCENT  = RGBColor(6, 182, 212)    # Cyan-500 — used surgically
C_TEXT_HI = RGBColor(255, 255, 255)  # Primary headings
C_TEXT_MID= RGBColor(148, 163, 184)  # Body copy — cool gray
C_TEXT_LOW= RGBColor(71, 85, 105)    # Muted labels / numbers
C_RULE    = RGBColor(30, 41, 59)     # Subtle divider lines

# Grid
W, H      = Inches(13.333), Inches(7.5)
ML, MR    = Inches(1.0), Inches(1.0)
MT        = Inches(0.7)
CW        = W - ML - MR              # Content width ≈ 11.33"

# Typography
FONT_HEAD = "Calibri Light"
FONT_BODY = "Calibri"


# ==============================================================================
# PRIMITIVE HELPERS — Thin abstraction layer over python-pptx
# ==============================================================================

def _bg(slide):
    """Full-bleed background. Added first so it sits at z-index 0."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_BG
    shape.line.fill.background()


def _rule(slide, left, top, width, height=Inches(0.012), color=C_RULE):
    """Hairline rectangle for dividers and accent anchors."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def _text(slide, left, top, width, height, text, size=Pt(14), bold=False,
          color=C_TEXT_MID, align=PP_ALIGN.LEFT, font=FONT_BODY, line=1.25):
    """Single-run textbox with zero internal padding."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = size
    r.font.bold = bold
    r.font.color.rgb = color
    return tb, tf


def _label(slide, text, top=MT):
    """Small uppercase section label with a short accent underline."""
    _text(slide, ML, top, CW, Inches(0.35), text.upper(),
          size=Pt(10), bold=True, color=C_ACCENT, font=FONT_BODY)
    _rule(slide, ML, top + Inches(0.32), Inches(1.4), height=Inches(0.015), color=C_ACCENT)


def _heading(slide, text, top, align=PP_ALIGN.LEFT):
    """Large display heading."""
    _text(slide, ML, top, CW, Inches(0.9), text,
          size=Pt(34), bold=True, color=C_TEXT_HI, align=align, font=FONT_HEAD, line=1.1)


# ==============================================================================
# SLIDE BUILDERS — Each slide gets a bespoke layout, no reused card template
# ==============================================================================

def slide_title(prs, blank):
    s = prs.slides.add_slide(blank)
    _bg(s)

    # Top hairline — full bleed, 1 px feel
    _rule(s, Inches(0), Inches(0), W, height=Inches(0.03), color=C_ACCENT)

    # Conference label
    _text(s, ML, Inches(2.0), CW, Inches(0.4),
          "GSMA MENA IGNITE HACKATHON 2026",
          size=Pt(12), bold=True, color=C_TEXT_LOW, align=PP_ALIGN.CENTER, font=FONT_BODY)

    # Massive title — typography IS the image
    tb = s.shapes.add_textbox(ML, Inches(2.5), CW, Inches(1.3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "AEGISTEL"
    r.font.name = FONT_HEAD
    r.font.size = Pt(80)
    r.font.bold = True
    r.font.color.rgb = C_TEXT_HI

    # Centered accent anchor
    _rule(s, Inches(5.9), Inches(3.75), Inches(1.533), height=Inches(0.025), color=C_ACCENT)

    # Subtitle
    _text(s, ML, Inches(4.0), CW, Inches(0.5),
          "Autonomous Telco-Aware AI Guard Engine",
          size=Pt(24), bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER, font=FONT_HEAD)

    # Descriptor — narrow measure for readability
    _text(s, ML + Inches(2.0), Inches(4.7), CW - Inches(4.0), Inches(0.8),
          "Orchestrating 6 GSMA CAMARA APIs for Real-Time 5G Network Intelligence",
          size=Pt(16), color=C_TEXT_MID, align=PP_ALIGN.CENTER, line=1.35)


def slide_problem_solution(prs, blank):
    s = prs.slides.add_slide(blank)
    _bg(s)
    _label(s, "01  Problem & Solution", top=Inches(0.6))
    _heading(s, "The Challenge & Our Response", top=Inches(1.15))

    # Vertical divider — subtle, architectural
    div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(6.666), Inches(2.0), Inches(0.008), Inches(4.6))
    div.fill.solid()
    div.fill.fore_color.rgb = C_RULE
    div.line.fill.background()

    col_w = Inches(5.0)
    left_x, right_x = ML, Inches(7.15)

    # --- LEFT: Problem ---
    _rule(s, left_x, Inches(2.0), Inches(0.7), height=Inches(0.03), color=C_ACCENT)
    _text(s, left_x, Inches(2.05), col_w, Inches(0.35),
          "THE CHALLENGE", size=Pt(13), bold=True, color=C_ACCENT)

    problems = [
        ("Vulnerable Authentication",
         "Legacy reliance on SMS OTPs exposes users to SIM-swap fraud and intercept attacks."),
        ("Siloed Network Signals",
         "Modern 5G features remain disconnected from application logic."),
        ("Manual Emergency Priority",
         "Public safety teams lack automated access to 5G priority slicing."),
        ("Developer Complexity",
         "Enterprise engineers find 3GPP network standards overly complex to program."),
    ]
    y = Inches(2.55)
    for title, desc in problems:
        _text(s, left_x, y, col_w, Inches(0.3), title,
              size=Pt(15), bold=True, color=C_TEXT_HI)
        _text(s, left_x, y + Inches(0.28), col_w, Inches(0.7), desc,
              size=Pt(12), color=C_TEXT_MID, line=1.35)
        y += Inches(1.15)

    # --- RIGHT: Solution ---
    _rule(s, right_x, Inches(2.0), Inches(0.7), height=Inches(0.03), color=C_ACCENT)
    _text(s, right_x, Inches(2.05), col_w, Inches(0.35),
          "THE SOLUTION", size=Pt(13), bold=True, color=C_ACCENT)

    solutions = [
        ("Autonomous AI Orchestration",
         "Groq LLM dynamically chains 6 CAMARA network APIs in real time."),
        ("Zero-Trust Verification",
         "Replaces OTPs with silent carrier identity and instant SIM-swap checks."),
        ("Proactive 5G Quality-on-Demand",
         "Automatically provisions high-priority 5G slices for emergency teams."),
        ("Unified Telco Guard Layer",
         "Simplifies complex 3GPP interfaces into zero-code autonomous triggers."),
    ]
    y = Inches(2.55)
    for title, desc in solutions:
        _text(s, right_x, y, col_w, Inches(0.3), title,
              size=Pt(15), bold=True, color=C_TEXT_HI)
        _text(s, right_x, y + Inches(0.28), col_w, Inches(0.7), desc,
              size=Pt(12), color=C_TEXT_MID, line=1.35)
        y += Inches(1.15)


def slide_api_grid(prs, blank):
    s = prs.slides.add_slide(blank)
    _bg(s)
    _label(s, "02  Proposed Solution", top=Inches(0.6))
    _heading(s, "Six CAMARA APIs", top=Inches(1.15))

    apis = [
        ("SIM Swap Detection", "Flags recent SIM changes to prevent account takeover."),
        ("Location Verification", "3GPP-level geofencing without device GPS dependency."),
        ("Quality on Demand", "Allocates priority 5QI 5G network slices live."),
        ("Number Verification", "Silent carrier identity matching without SMS OTPs."),
        ("Congestion Insights", "Monitors cell density for smart city safety."),
        ("Device Reachability", "Checks connectivity before SLA adjustments."),
    ]

    card_w = Inches(3.4)
    card_h = Inches(2.1)
    gap_x, gap_y = Inches(0.45), Inches(0.35)
    start_y = Inches(2.05)

    for idx, (title, desc) in enumerate(apis):
        col, row = idx % 3, idx // 3
        x = ML + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        # Large muted number — acts as the only "graphic"
        _text(s, x, y, Inches(0.6), Inches(0.5), f"0{idx+1}",
              size=Pt(32), bold=True, color=C_TEXT_LOW, font=FONT_HEAD)

        # Title and description indented from the number
        _text(s, x + Inches(0.55), y + Inches(0.08), card_w - Inches(0.7), Inches(0.35),
              title, size=Pt(16), bold=True, color=C_TEXT_HI)
        _text(s, x + Inches(0.55), y + Inches(0.42), card_w - Inches(0.7), Inches(0.8),
              desc, size=Pt(12), color=C_TEXT_MID, line=1.3)

        # Bottom hairline separator per item
        _rule(s, x, y + card_h - Inches(0.06), card_w, height=Inches(0.008), color=C_RULE)


def slide_ai_flow(prs, blank):
    s = prs.slides.add_slide(blank)
    _bg(s)
    _label(s, "03  AI Agent Design", top=Inches(0.6))
    _heading(s, "Autonomous Decision Flow", top=Inches(1.15))

    steps = [
        ("Intent Parsing", "Receives event payloads and extracts operational intent."),
        ("Tool Selection", "Groq Llama-3.3-70B dynamically selects the optimal CAMARA APIs."),
        ("Parallel Execution", "Fires simultaneous verification calls to Nokia NaC."),
        ("Telemetry Fusion", "Merges risk indicators into a structured composite score."),
        ("Policy Decision", "Emits ALLOW, BLOCK, or ESCALATE over WebSockets."),
    ]

    y0 = Inches(2.1)
    step_h = Inches(1.0)

    for idx, (title, desc) in enumerate(steps):
        y = y0 + idx * step_h
        cx = ML + Inches(0.15)   # circle center x offset

        # Number circle — filled dark, outlined accent
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, y, Inches(0.32), Inches(0.32))
        circ.fill.solid()
        circ.fill.fore_color.rgb = RGBColor(12, 30, 50)
        circ.line.color.rgb = C_ACCENT
        circ.line.width = Pt(1.5)

        # Number centered in circle
        tb = s.shapes.add_textbox(cx, y + Inches(0.02), Inches(0.32), Inches(0.28))
        tf = tb.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = str(idx + 1)
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = C_ACCENT

        # Connector line between circles
        if idx < len(steps) - 1:
            conn = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      cx + Inches(0.15), y + Inches(0.32),
                                      Inches(0.02), Inches(0.68))
            conn.fill.solid()
            conn.fill.fore_color.rgb = C_RULE
            conn.line.fill.background()

        # Text block
        tx = ML + Inches(0.75)
        _text(s, tx, y + Inches(0.02), CW - Inches(1.2), Inches(0.3),
              title, size=Pt(16), bold=True, color=C_TEXT_HI)
        _text(s, tx, y + Inches(0.32), CW - Inches(1.2), Inches(0.4),
              desc, size=Pt(13), color=C_TEXT_MID, line=1.3)


def slide_architecture(prs, blank):
    s = prs.slides.add_slide(blank)
    _bg(s)
    _label(s, "04  Technical Architecture", top=Inches(0.6))
    _heading(s, "Where 5G Fits", top=Inches(1.15))

    items = [
        ("5G Network Exposure Function (NEF)",
         "Northbound 3GPP interface translating AI agent decisions into core network instructions."),
        ("5G Network Slicing & 5QI Customization",
         "Dynamic binding of active traffic to QoS Identifiers via the Policy Control Function."),
        ("5G Core Control Plane Integration",
         "Direct API integration with UDM and GMLC for identity and location services."),
        ("Multi-Access Edge Computing (MEC)",
         "Directs high-density congestion traffic and emergency video streams to local edge nodes."),
    ]

    y = Inches(2.0)
    for title, desc in items:
        # Left accent bar — strong vertical anchor
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, ML, y, Inches(0.04), Inches(0.9))
        bar.fill.solid()
        bar.fill.fore_color.rgb = C_ACCENT
        bar.line.fill.background()

        _text(s, ML + Inches(0.18), y + Inches(0.06), CW - Inches(0.6), Inches(0.35),
              title, size=Pt(17), bold=True, color=C_TEXT_HI)
        _text(s, ML + Inches(0.18), y + Inches(0.38), CW - Inches(0.6), Inches(0.5),
              desc, size=Pt(13), color=C_TEXT_MID, line=1.3)

        y += Inches(1.35)


def slide_business(prs, blank):
    s = prs.slides.add_slide(blank)
    _bg(s)
    _label(s, "05  Business Model", top=Inches(0.6))
    _heading(s, "Commercial Value", top=Inches(1.15))

    col_w = Inches(5.2)
    left_x, right_x = ML, Inches(7.1)

    # --- LEFT ---
    _rule(s, left_x, Inches(1.95), Inches(0.9), height=Inches(0.025), color=C_ACCENT)
    _text(s, left_x, Inches(2.0), col_w, Inches(0.35),
          "MONETIZATION", size=Pt(12), bold=True, color=C_ACCENT)

    monetization = [
        "Telco Revenue Sharing — B2B API usage fees per CAMARA invocation.",
        "SaaS Enterprise Tier — Monthly subscription for banks and smart city platforms.",
        "Premium 5G Slice Surcharges — Pay-per-use for emergency QoD allocations.",
    ]
    y = Inches(2.55)
    for item in monetization:
        # Custom bullet: small square
        sq = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, y + Inches(0.09),
                                Inches(0.06), Inches(0.06))
        sq.fill.solid()
        sq.fill.fore_color.rgb = C_ACCENT
        sq.line.fill.background()

        _text(s, left_x + Inches(0.18), y, col_w - Inches(0.3), Inches(0.8),
              item, size=Pt(14), color=C_TEXT_MID, line=1.35)
        y += Inches(0.95)

    # --- RIGHT ---
    _rule(s, right_x, Inches(1.95), Inches(0.9), height=Inches(0.025), color=C_ACCENT)
    _text(s, right_x, Inches(2.0), col_w, Inches(0.35),
          "IMPACT", size=Pt(12), bold=True, color=C_ACCENT)

    impact = [
        "Fraud Reduction — Over 90% reduction in SIM-swap account takeover losses.",
        "Zero Friction UX — Eliminates manual OTPs, increasing completion by 25%.",
        "Network Monetization — Unlocks new enterprise revenues from idle 5G assets.",
    ]
    y = Inches(2.55)
    for item in impact:
        sq = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, y + Inches(0.09),
                                Inches(0.06), Inches(0.06))
        sq.fill.solid()
        sq.fill.fore_color.rgb = C_ACCENT
        sq.line.fill.background()

        _text(s, right_x + Inches(0.18), y, col_w - Inches(0.3), Inches(0.8),
              item, size=Pt(14), color=C_TEXT_MID, line=1.35)
        y += Inches(0.95)


def slide_demo(prs, blank):
    s = prs.slides.add_slide(blank)
    _bg(s)
    _label(s, "06  Virtual Demo", top=Inches(0.6))
    _heading(s, "Prototype & Resources", top=Inches(1.15))

    items = [
        ("Demo Video", "3-minute screen recording of live API calls, WebSocket traces, and dashboard."),
        ("Repository", "github.com/nullstack/aegistel-mena-ignite"),
        ("Frontend", "Vite + React + Tailwind dashboard rendering agent execution steps."),
        ("Backend", "FastAPI gateway with OpenAPI docs and WebSocket streaming on /ws/agent."),
    ]

    y = Inches(2.2)
    for num, (title, desc) in enumerate(items, 1):
        # Number
        _text(s, ML, y, Inches(0.5), Inches(0.3), f"0{num}",
              size=Pt(14), bold=True, color=C_ACCENT)

        # Title
        _text(s, ML + Inches(0.5), y, CW - Inches(1.5), Inches(0.3),
              title, size=Pt(16), bold=True, color=C_TEXT_HI)

        # Description
        _text(s, ML + Inches(0.5), y + Inches(0.3), CW - Inches(1.5), Inches(0.4),
              desc, size=Pt(13), color=C_TEXT_MID)

        # Separator
        _rule(s, ML + Inches(0.5), y + Inches(0.75), CW - Inches(2.0),
              height=Inches(0.006), color=C_RULE)

        y += Inches(1.15)


def slide_team(prs, blank):
    s = prs.slides.add_slide(blank)
    _bg(s)
    _label(s, "07  Team", top=Inches(0.6))

    # Large name — left aligned, commanding presence
    tb = s.shapes.add_textbox(ML, Inches(2.1), CW, Inches(1.0))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Yahia Abdeldjalil"
    r.font.name = FONT_HEAD
    r.font.size = Pt(52)
    r.font.bold = True
    r.font.color.rgb = C_TEXT_HI

    # Accent rule
    _rule(s, ML, Inches(2.85), Inches(2.2), height=Inches(0.02), color=C_ACCENT)

    # Role
    _text(s, ML, Inches(3.0), CW, Inches(0.4),
          "AI Systems Engineer & Full-Stack Developer",
          size=Pt(18), bold=True, color=C_ACCENT, font=FONT_HEAD)

    # Bio lines — narrow measure for elegance
    bios = [
        "Architected the Groq-powered multi-agent orchestration engine with 6 CAMARA APIs.",
        "Integrated Nokia NaC and developed the real-time React dashboard.",
        "Specialization: AI / LLM Tool Orchestration & 5G Programmable Core Interfaces.",
    ]
    y = Inches(3.7)
    for bio in bios:
        _text(s, ML, y, CW - Inches(3.0), Inches(0.5), bio,
              size=Pt(14), color=C_TEXT_MID, line=1.35)
        y += Inches(0.6)


# ==============================================================================
# MAIN
# ==============================================================================

def create_presentation():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]  # Blank

    slide_title(prs, blank)
    slide_problem_solution(prs, blank)
    slide_api_grid(prs, blank)
    slide_ai_flow(prs, blank)
    slide_architecture(prs, blank)
    slide_business(prs, blank)
    slide_demo(prs, blank)
    slide_team(prs, blank)

    out = "aegistel_keynote.pptx"
    prs.save(out)
    print(f"Generated {out}")


if __name__ == "__main__":
    create_presentation()